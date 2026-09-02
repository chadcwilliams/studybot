"""
Parse course materials (.pptx, .pdf, .docx) in data/materials/, split them
into overlapping text chunks, embed them locally, and persist them to a
local vector store on disk (see store.py).

Run this once after adding/updating files:

    python -m studybot.ingest

Re-run it any time you add new slides or update the syllabus — it rebuilds
the index from scratch each time, so it's always safe to re-run.
"""

from __future__ import annotations

import base64
import hashlib
import json
import re
import sys
import time
import xml.etree.ElementTree as ET
from dataclasses import dataclass
from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document as DocxDocument

# dwml is a nice light dependency (pure Python + stdlib xml), so it's a
# normal top-level import rather than the lazy pattern used for the heavy
# ML libraries elsewhere in this file.
from dwml import omml

from studybot.config import settings
from studybot.store import StoredChunk, VectorStore

# fastembed is imported lazily inside run() below, so the text-extraction
# functions above can be unit-tested or reused without downloading models.


@dataclass
class Chunk:
    text: str
    source: str  # filename
    location: str  # e.g. "slide 4" or "page 2"


# --------------------------------------------------------------------------
# Text extraction
# --------------------------------------------------------------------------

def extract_pptx(path: Path) -> list[tuple[str, str]]:
    """Returns list of (text, location) per slide, including speaker notes."""
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(c.text.strip() for c in row.cells)
                    if row_text.strip(" |"):
                        parts.append(row_text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            parts.append("Speaker notes: " + slide.notes_slide.notes_text_frame.text.strip())
        if parts:
            out.append(("\n".join(parts), f"slide {i}"))
    return out


def extract_pdf(path: Path) -> list[tuple[str, str]]:
    reader = PdfReader(str(path))
    out = []
    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            out.append((text, f"page {i}"))
    return out


def _text_with_math(docx_element) -> str:
    """Combines a paragraph/cell's plain text with any embedded Word
    equation objects (OMML). python-docx's .text property silently skips
    over these entirely — unlike PowerPoint's embedded pictures (see
    describe_pptx_images below), Word's equation editor produces real
    structured math markup, not an image, so it converts to LaTeX directly
    and exactly: no AI model, no cost, no ambiguity.

    Note: this appends any equations found after the plain text rather than
    interleaving them in exact document order. That matches the common
    real-world pattern (a label like "SD =" followed by the equation) but
    won't perfectly reconstruct a cell with text before AND after an
    equation — good enough for formula sheets and tables without being a
    full XML-order reconstruction.
    """
    plain = docx_element.text.strip()

    xml = None
    if hasattr(docx_element, "_tc"):
        xml = docx_element._tc.xml
    elif hasattr(docx_element, "_p"):
        xml = docx_element._p.xml
    if xml is None:
        return plain

    try:
        root = ET.fromstring(xml)
    except ET.ParseError:
        return plain

    latex_parts = []
    for elem in root.iter(f"{omml.OMML_NS}oMath"):
        try:
            latex = omml.oMath2Latex(elem).latex
        except Exception:  # noqa: BLE001 - a malformed equation shouldn't break extraction
            continue
        if latex and latex.strip():
            latex_parts.append(latex.strip())

    if not latex_parts:
        return plain

    math_text = " ".join(f"${l}$" for l in latex_parts)
    return f"{plain} {math_text}".strip()


def extract_docx(path: Path) -> list[tuple[str, str]]:
    doc = DocxDocument(path)
    out = []

    # Split the narrative text into sections at heading boundaries, instead
    # of joining every paragraph into one giant blob. A single ~19,000-char
    # blob gets sliced later purely by raw character count (chunk_text),
    # with no regard for where one policy ends and another begins — the same
    # problem we already fixed for tables, just showing up in prose instead.
    #
    # Word's "Heading" styles are the obvious signal to split on, but some
    # documents apply a heading style to body paragraphs too (not just their
    # titles) — treating every Heading-styled paragraph as a new, empty
    # section wipes that content out rather than just misplacing it. So a
    # paragraph only counts as a real heading if it also looks like one:
    # short, and not ending in sentence-final punctuation. A body paragraph
    # that happens to carry a Heading style almost never satisfies both.
    def _is_real_heading(text: str, style_name: str) -> bool:
        if not style_name.startswith("Heading"):
            return False
        return len(text) <= 80 and not text.rstrip().endswith((".", "!", "?"))

    sections: list[tuple[str, list[str]]] = []
    current_heading = "Introduction"
    current_paras: list[str] = []

    def flush() -> None:
        if current_paras:
            sections.append((current_heading, list(current_paras)))

    for p in doc.paragraphs:
        text = _text_with_math(p)
        if not text:
            continue
        if _is_real_heading(text, p.style.name):
            flush()
            current_heading = text
            current_paras = []
        else:
            current_paras.append(text)
    flush()

    for heading, paras in sections:
        body = "\n".join(paras)
        if not body.strip():
            continue
        section_text = body if heading == "Introduction" else f"{heading}\n{body}"
        out.append((section_text, f"document text — {heading[:60]}"))

    # Tables (grading breakdowns, exam schedules, etc.) live in doc.tables,
    # separate from doc.paragraphs. We chunk them row-by-row — each row
    # paired with its header row — rather than joining the whole table into
    # one blob. A long table joined into one string can get sliced by the
    # generic character-based chunker at an arbitrary point, severing a fact
    # (e.g. "Exam 1" from its date "Oct 15") across two separate chunks.
    # Row-level chunks are small enough to never get split, and keeping the
    # header attached preserves what each column means.
    for ti, table in enumerate(doc.tables, start=1):
        rows = [[_text_with_math(cell) for cell in row.cells] for row in table.rows]
        rows = [r for r in rows if any(c for c in r)]
        if not rows:
            continue

        if len(rows) == 1:
            # A single-row table is usually several side-by-side categories
            # (e.g. "Exams: 45%" | "Article Summary: 45%" | "Participation:
            # 2%") stuffed into one row of cells. Joining them with " | "
            # let facts from one category bleed into another, so each
            # column gets its own chunk — good for narrow questions about
            # one category.
            non_empty_cells = [c.strip() for c in rows[0] if c.strip()]
            for ci, cell_text in enumerate(non_empty_cells, start=1):
                out.append((cell_text, f"table {ti}, column {ci}"))

            # But a BROAD question ("what's the whole grade breakdown?")
            # needs all of those columns retrieved together — a much
            # harder bar than matching one chunk. Add one more chunk with
            # every column present, separated by blank lines (not " | ")
            # so categories stay visually distinct and don't get
            # cross-attributed, while still being answerable from a
            # single retrieved chunk.
            if len(non_empty_cells) > 1:
                out.append(("\n\n".join(non_empty_cells), f"table {ti} (all columns)"))
            continue

        header = rows[0]
        # A row-0 cell containing "$" (our own LaTeX marker) or "=" is a
        # reliable sign that row 0 is actual data (e.g. a formula), not a
        # column-label header — a real header describes what KIND of data
        # is in a column ("Date", "Topic"), it doesn't contain an instance
        # of that data itself. Treating a headerless table's first row as a
        # header wrongly prepends it onto every subsequent row.
        has_real_header = not any("$" in c or "=" in c for c in header)

        if not has_real_header:
            # No column labels to attach — every row (including row 0) is
            # its own standalone item (e.g. one row per named equation).
            # Just join each row's own non-empty cells.
            for ri, row in enumerate(rows, start=1):
                cells = [c.strip() for c in row if c.strip()]
                if cells:
                    out.append((": ".join(cells), f"table {ti}, row {ri}"))
            continue

        date_col = next((i for i, h in enumerate(header) if "date" in h.lower()), None)

        for ri, row in enumerate(rows[1:], start=2):
            # Build a natural "label: value" description using only the
            # non-empty cells in THIS row, instead of always repeating the
            # full header line. That header repetition was the problem: with
            # it, all 27 rows of a schedule table carry faint similarity to
            # any query mentioning a column name (like "due"), even rows
            # that have nothing due — drowning out the rows that actually
            # matter. Omitting empty cells means a row only "sounds like"
            # what it actually contains.
            pairs = [f"{header[i]}: {row[i]}" for i in range(len(row)) if row[i].strip()]
            if not pairs:
                continue
            out.append((", ".join(pairs), f"table {ti}, row {ri}"))

        # A long schedule table with many near-identical rows (same columns,
        # slightly different dates/topics) is hard for embedding search to
        # pick out of — a row like "Summary 1" differs from its 27 neighbors
        # by one word, which is too weak a signal to reliably rank near the
        # top for a question like "when are the assignments due". If a
        # column's header suggests it marks deadlines, build one consolidated
        # chunk listing every non-empty entry in that column as a plain
        # sentence — natural phrasing matches how students actually ask,
        # and one dense, unambiguous chunk beats hoping semantic search finds
        # all the scattered needle rows on its own.
        due_col = next(
            (i for i, h in enumerate(header) if any(kw in h.lower() for kw in ("due", "deadline"))),
            None,
        )
        if due_col is not None:
            sentences = []
            for row in rows[1:]:
                value = row[due_col].strip()
                if not value:
                    continue
                date_val = row[date_col].strip() if date_col is not None else ""
                other_pairs = [
                    f"{header[i]}: {row[i]}"
                    for i in range(len(row))
                    if i not in (due_col, date_col) and row[i].strip()
                ]
                other_context = ", ".join(other_pairs)
                if date_val:
                    sentence = f"{value} is due on {date_val}"
                    if other_context:
                        sentence += f" ({other_context})"
                    sentence += "."
                else:
                    sentence = f"{value} — {other_context}." if other_context else f"{value}."
                sentences.append(sentence)
            if sentences:
                due_header = header[due_col].strip()
                summary = f"Full list of everything under \"{due_header}\":\n" + "\n".join(sentences)
                out.append((summary, f"table {ti} — {due_header} (summary)"))

    return out


EXTRACTORS = {
    ".pptx": extract_pptx,
    ".pdf": extract_pdf,
    ".docx": extract_docx,
}


# --------------------------------------------------------------------------
# Chunking
# --------------------------------------------------------------------------

def chunk_text(text: str, size: int, overlap: int) -> list[str]:
    """Simple character-based sliding-window chunker."""
    text = text.strip()
    if len(text) <= size:
        return [text] if text else []
    chunks = []
    start = 0
    while start < len(text):
        end = start + size
        chunks.append(text[start:end])
        if end >= len(text):
            break
        start = end - overlap
    return chunks


def build_chunks(materials_dir: Path, chunk_size: int, chunk_overlap: int) -> list[Chunk]:
    chunks: list[Chunk] = []
    files = [f for f in materials_dir.rglob("*") if f.suffix.lower() in EXTRACTORS]

    if not files:
        print(f"No supported files found in {materials_dir}. "
              f"Add .pptx / .pdf / .docx files there and re-run.")
        return chunks

    for f in files:
        extractor = EXTRACTORS[f.suffix.lower()]
        try:
            sections = extractor(f)
        except Exception as e:  # noqa: BLE001 - keep ingesting other files
            print(f"  ! Failed to parse {f.name}: {e}", file=sys.stderr)
            continue

        for text, location in sections:
            for piece in chunk_text(text, chunk_size, chunk_overlap):
                if piece.strip():
                    chunks.append(Chunk(text=piece, source=f.name, location=location))

        print(f"  parsed {f.name}: {len(sections)} section(s)")

    return chunks


# --------------------------------------------------------------------------
# Image description (equations/diagrams embedded as pictures — invisible to
# plain text extraction, so a vision model transcribes/describes them once
# at ingestion time, not per student question).
# --------------------------------------------------------------------------

_IMAGE_DESCRIBE_PROMPT = (
    "This image is from a course slide. If it shows a mathematical equation, "
    "formula, or notation, transcribe it exactly and completely using LaTeX "
    "($ for inline, $$ for standalone). If it's a diagram, chart, graph, or "
    "labeled figure, describe its content and every label precisely enough "
    "that someone could answer questions about it without seeing the image. "
    "If it's purely decorative (a logo, background texture, or photo with no "
    "course-relevant information), respond with exactly: DECORATIVE_IMAGE_SKIP"
)


def _load_image_cache(path: Path) -> dict[str, str]:
    if not path.exists():
        return {}
    try:
        with open(path, encoding="utf-8") as f:
            return json.load(f)
    except (OSError, json.JSONDecodeError):
        return {}


def _save_image_cache(path: Path, cache: dict[str, str]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(cache, f, indent=2)


class _VisionQuotaExhausted(Exception):
    """Raised when a failure looks like a hard daily/token quota (not a
    transient per-minute rate limit) -- signals the caller to stop trying
    further images this run rather than failing on every remaining one
    individually, since they'll all fail identically until the quota resets.
    """


_QUOTA_DETAIL_PATTERN = re.compile(
    r"Limit (?P<limit>[\d,]+), Used (?P<used>[\d,]+), Requested (?P<requested>[\d,]+)\.?"
    r".*?try again in (?P<retry>[^.]+?\.?\d*s)",
    re.IGNORECASE | re.DOTALL,
)


def _summarize_quota_error(message: str) -> str:
    """Pulls the useful numbers out of Groq's error text and formats them
    readably. Falls back to the raw message if the format doesn't match —
    Groq could change their wording, and a failed parse shouldn't hide the
    underlying error.
    """
    match = _QUOTA_DETAIL_PATTERN.search(message)
    if not match:
        return message
    limit = int(match.group("limit").replace(",", ""))
    used = int(match.group("used").replace(",", ""))
    remaining = max(0, limit - used)
    retry = match.group("retry").strip()
    pct = (used / limit * 100) if limit else 0
    return (
        f"{used:,} / {limit:,} tokens used today ({pct:.1f}%), "
        f"{remaining:,} remaining. Retry in {retry}."
    )


def _describe_image(client, model: str, blob: bytes, content_type: str) -> str | None:
    """Returns a description, "" if genuinely decorative (this gets cached --
    correctly never retried), or None if the call failed for an ordinary
    reason (NOT cached, so it gets retried on the next ingest instead of
    being permanently mistaken for "decorative"). Raises
    _VisionQuotaExhausted if the failure specifically looks like the daily
    token quota being exhausted, rather than returning None for that case --
    the caller needs to know to stop entirely, not just skip this one image.
    """
    data_url = f"data:{content_type};base64,{base64.b64encode(blob).decode()}"
    try:
        response = client.chat.completions.create(
            model=model,
            temperature=0.1,
            max_tokens=500,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": _IMAGE_DESCRIBE_PROMPT},
                    {"type": "image_url", "image_url": {"url": data_url}},
                ],
            }],
        )
        text = response.choices[0].message.content.strip()
        return "" if text == "DECORATIVE_IMAGE_SKIP" else text
    except Exception as e:  # noqa: BLE001 - one bad image shouldn't abort ingestion
        message = str(e)
        if "rate_limit_exceeded" in message and ("per day" in message.lower() or "tpd" in message.lower()):
            raise _VisionQuotaExhausted(_summarize_quota_error(message)) from e
        print(f"    ! image description failed (will retry next ingest): {e}", file=sys.stderr)
        return None


def describe_pptx_images(materials_dir: Path) -> list[Chunk]:
    from groq import Groq

    if not settings.groq_api_key:
        print("  ! GROQ_API_KEY not set -- skipping image description.")
        return []

    client = Groq(api_key=settings.groq_api_key)
    cache = _load_image_cache(settings.image_cache_path)
    chunks: list[Chunk] = []

    for f in materials_dir.rglob("*.pptx"):
        try:
            prs = Presentation(f)
        except Exception as e:  # noqa: BLE001
            print(f"  ! Failed to open {f.name} for image scan: {e}", file=sys.stderr)
            continue

        for i, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                try:
                    image = shape.image
                except Exception:  # noqa: BLE001 - some picture shapes have no accessible image data
                    continue
                if len(image.blob) < settings.image_min_bytes:
                    continue  # almost certainly an icon or decorative dot

                image_hash = hashlib.sha256(image.blob).hexdigest()
                if image_hash in cache:
                    description = cache[image_hash]
                else:
                    print(f"  describing image: {f.name}, slide {i}...")
                    try:
                        description = _describe_image(client, settings.vision_model, image.blob, image.content_type)
                    except _VisionQuotaExhausted as e:
                        print(
                            "  ! Vision model's daily quota is exhausted. Stopping image "
                            "description for this run -- already-described images are saved "
                            "(saved incrementally, not just at the end).\n"
                            f"    {e}"
                        )
                        return chunks
                    time.sleep(2)  # pace requests under the vision model's own rate limit
                    if description is None:
                        continue  # failed call -- not cached, will retry next ingest
                    cache[image_hash] = description
                    # Save after EVERY new description, not just at the end -- if the
                    # process is interrupted (Ctrl+C, closed terminal, crash) partway
                    # through a long run, already-completed work isn't lost.
                    _save_image_cache(settings.image_cache_path, cache)

                if description:
                    chunks.append(Chunk(text=description, source=f.name, location=f"slide {i} -- image"))

    return chunks


# --------------------------------------------------------------------------
# Embedding + storage
# --------------------------------------------------------------------------

def run() -> None:
    from fastembed import TextEmbedding

    print(f"Reading materials from: {settings.materials_dir}")
    chunks = build_chunks(settings.materials_dir, settings.chunk_size, settings.chunk_overlap)

    if settings.describe_images:
        print("Describing embedded images (cached by content hash — only new "
              "images cost a vision-model call)...")
        chunks.extend(describe_pptx_images(settings.materials_dir))

    if not chunks:
        print("Nothing to index. Exiting.")
        return

    print(f"Built {len(chunks)} chunks. Loading embedding model "
          f"'{settings.embedding_model}' (first run downloads it)...")
    model = TextEmbedding(settings.embedding_model)

    print("Embedding chunks...")
    texts = [c.text for c in chunks]
    embeddings = [e.tolist() for e in model.embed(texts)]

    store = VectorStore(settings.index_dir)
    stored_chunks = [StoredChunk(text=c.text, source=c.source, location=c.location) for c in chunks]
    store.save(embeddings, stored_chunks)

    print(f"Done. Indexed {len(chunks)} chunks at {settings.index_dir}")


if __name__ == "__main__":
    run()
